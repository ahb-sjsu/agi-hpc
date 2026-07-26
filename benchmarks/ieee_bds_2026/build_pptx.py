# Rebuild SecureAI_Bond_2026_slides.tex as a native, editable PPTX.
# Restyled 2026-07-25: flat title treatment, ColorBrewer-derived palette
# (Blues #2171B5 / Dark2 orange #D95F02 / Dark2 teal #1B9E77 — validated
# for CVD separation and 3:1 surface contrast with the dataviz checker).
# Peer-review fixes applied; deck order mirrors the .tex:
#   1 title, 2 problem, 3 approach, 4 finding-1, 5 finding-2, 6 validation,
#   7 exploit, 8 kernel, 9 defense, 10 bifactor (native chart),
#   11 implications, 12 philosophy, 13 takeaways, 14 thanks,
#   15-16 backup (grounding table, discovery loop).
# NOTES has exactly one entry per slide, same order — keep them in sync.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
import copy

ASSETS = r"C:\source\agi-hpc\benchmarks\ieee_bds_2026\revision\assets"
OUT = r"C:\source\agi-hpc\benchmarks\ieee_bds_2026\SecureAI_Bond_2026_slides.pptx"

# ColorBrewer-derived palette (validated). Variable names kept generic so
# every add_runs style keeps working.
ATLASBLUE = RGBColor(0x21, 0x71, 0xB5)  # Blues — structure / general factor
ERISGOLD = RGBColor(0xD9, 0x5F, 0x02)   # Dark2 orange — accent / specific
CBTEAL = RGBColor(0x1B, 0x9E, 0x77)     # Dark2 teal — mixed / defended
INK = RGBColor(0x1A, 0x23, 0x40)
BLOCKBG = RGBColor(0xE4, 0xED, 0xF7)    # cbblue ~8% on white
GRAY = RGBColor(0x63, 0x63, 0x63)
BLACK = RGBColor(0x20, 0x20, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ---- helpers ---------------------------------------------------------------
# a "segment" is (text, style); style: None, 'hi' (blue bold), 'gold' (orange
# bold), 'em' (italic), 'bold', 'mono' (Consolas)


def add_runs(p, segments, size, base_color=BLACK):
    for text, style in segments:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = FONT
        f.size = Pt(size)
        f.color.rgb = base_color
        if style == "hi":
            f.bold = True
            f.color.rgb = ATLASBLUE
        elif style == "gold":
            f.bold = True
            f.color.rgb = ERISGOLD
        elif style == "em":
            f.italic = True
        elif style == "bold":
            f.bold = True
        elif style == "mono":
            f.name = "Consolas"


def textbox(slide, x, y, w, h):
    shp = slide.shapes.add_textbox(x, y, w, h)
    tf = shp.text_frame
    tf.word_wrap = True
    return tf


def para(tf, segments, size=16, bullet=False, indent=0, space_after=6,
         align=PP_ALIGN.LEFT, color=BLACK, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if bullet:
        segments = [("\u25aa  ", "hi")] + list(segments)
    add_runs(p, segments, size, base_color=color)
    return p


def title_bar(slide, text, num=None):
    # flat title: ink bold text over a thin blue rule (no filled banner)
    tf = textbox(slide, Inches(0.45), Inches(0.16), SW - Inches(0.9), Inches(0.52))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = INK
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.74), SW - Inches(0.9), Pt(1.6))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ATLASBLUE
    rule.line.fill.background()
    if num:
        ftf = textbox(slide, SW - Inches(0.9), SH - Inches(0.42), Inches(0.7), Inches(0.3))
        para(ftf, [(str(num), None)], size=10, align=PP_ALIGN.RIGHT, color=GRAY, first=True)


def block(slide, x, y, w, h, title, body_lines, body_size=15):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = BLOCKBG
    shp.line.color.rgb = ATLASBLUE
    shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    add_runs(p, [(title, "hi")], body_size + 1)
    p.space_after = Pt(6)
    for seg in body_lines:
        para(tf, seg["segments"], size=body_size, bullet=seg.get("bullet", False),
             space_after=seg.get("space_after", 4))
    return shp


def table(slide, x, y, w, rows_data, col_widths, header=True, size=14,
          row_h=0.32):
    nrows, ncols = len(rows_data), len(rows_data[0])
    shp = slide.shapes.add_table(nrows, ncols, x, y, w, Inches(row_h * nrows))
    tbl = shp.table
    tbl.first_row = header
    tbl.horz_banding = False
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)
    for ri, row in enumerate(rows_data):
        for ci, cell_segs in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ATLASBLUE if (header and ri == 0) else WHITE
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            base = WHITE if (header and ri == 0) else BLACK
            segs = cell_segs if isinstance(cell_segs, list) else [(cell_segs, "bold" if header and ri == 0 else None)]
            add_runs(p, segs, size, base_color=base)
    return shp


def footnote(slide, segments, y=None, size=11):
    y = y if y is not None else SH - Inches(0.62)
    tf = textbox(slide, Inches(0.55), y, SW - Inches(1.1), Inches(0.5))
    para(tf, segments, size=size, color=GRAY, first=True)


def new_slide(title=None, num=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        title_bar(s, title, num)
    return s


# ---- diagram helpers (native shapes for the four concept figures) ----------

def fbox(slide, x, y, w, h, segments, border=ATLASBLUE, fill=None, size=10,
         dashed=False, bold_title=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.10
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill if fill is not None else WHITE
    shp.line.color.rgb = border
    shp.line.width = Pt(1.2)
    if dashed:
        ln = shp.line._get_or_add_ln()
        d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(d)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_runs(p, segments, size)
    return shp


def harrow(slide, x, y, w, color=GRAY, h=Inches(0.16)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.adjustments[0] = 0.55
    shp.adjustments[1] = 0.55
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def varrow(slide, x, y, h, color=GRAY, w=Inches(0.16), up=False):
    kind = MSO_SHAPE.UP_ARROW if up else MSO_SHAPE.DOWN_ARROW
    shp = slide.shapes.add_shape(kind, x, y, w, h)
    shp.adjustments[0] = 0.55
    shp.adjustments[1] = 0.55
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def dot(slide, x, y, color, r=Inches(0.055)):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - r, y - r, 2 * r, 2 * r)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def larrow(slide, x, y, w, color, h=Inches(0.10)):
    shp = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, x, y, w, h)
    shp.adjustments[0] = 0.5
    shp.adjustments[1] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


# ---- 1. title slide --------------------------------------------------------
s = prs.slides.add_slide(BLANK)
s.shapes.add_picture(ASSETS + r"\atlas_mark_light.png", Inches(5.92), Inches(0.55), height=Inches(1.5))
tf = textbox(s, Inches(1.2), Inches(2.3), Inches(10.93), Inches(3.9))
para(tf, [("Selective Invariance Violations in\nLarge Language Model Moral Judgment", "bold")],
     size=30, align=PP_ALIGN.CENTER, color=INK, first=True, space_after=10)
para(tf, [("A Geometric Framework for Behavioral Manipulation Detection", None)],
     size=19, align=PP_ALIGN.CENTER, color=ATLASBLUE, space_after=22)
para(tf, [("Andrew H. Bond", "bold")], size=17, align=PP_ALIGN.CENTER, color=INK, space_after=2)
para(tf, [("Dept. of Computer Engineering, San Jose State University", None)],
     size=14, align=PP_ALIGN.CENTER, color=GRAY, space_after=16)
para(tf, [("IEEE BigDataService 2026 — Special Track on Secure AI", "em")],
     size=14, align=PP_ALIGN.CENTER, color=INK)

# ---- 2. the problem ---------------------------------------------------------
s = new_slide("The problem: LLMs as decision services", 2)
tf = textbox(s, Inches(0.55), Inches(1.0), Inches(7.2), Inches(3.2))
para(tf, [("LLMs increasingly gate ", None), ("safety-critical decisions", "hi"),
          (" — content moderation, clinical triage, legal analysis.", None)],
     size=15, bullet=True, first=True, space_after=10)
para(tf, [("A secure evaluator should judge the ", None), ("facts", "hi"),
          (", not their ", None), ("presentation", "hi"), (".", None)],
     size=15, bullet=True, space_after=10)
para(tf, [("Behavioral manipulation vulnerability", "gold"),
          (": the attacker changes morally ", None), ("irrelevant", "em"),
          (" surface features — framing, tone, sensory detail — never the underlying situation.", None)],
     size=15, bullet=True)
# right-hand concept diagram: same facts -> two wordings -> opposite verdicts
LIGHT = RGBColor(0xF3, 0xF3, 0xF3)
fbox(s, Inches(9.15), Inches(0.95), Inches(2.6), Inches(0.45),
     [("the same facts", "bold")], border=INK, fill=LIGHT, size=11)
varrow(s, Inches(9.35), Inches(1.46), Inches(0.32))
varrow(s, Inches(11.55), Inches(1.46), Inches(0.32))
fbox(s, Inches(8.25), Inches(1.85), Inches(2.15), Inches(0.5),
     [("stated plainly", None)], border=GRAY, size=10)
fbox(s, Inches(10.65), Inches(1.85), Inches(2.15), Inches(0.5),
     [("euphemistic rewording", None)], border=GRAY, size=10)
varrow(s, Inches(9.22), Inches(2.42), Inches(0.32))
varrow(s, Inches(11.62), Inches(2.42), Inches(0.32))
fbox(s, Inches(8.75), Inches(2.82), Inches(1.05), Inches(0.42),
     [("FLAG", "hi")], border=ATLASBLUE, size=11)
fbox(s, Inches(11.15), Inches(2.82), Inches(1.05), Inches(0.42),
     [("PASS", "gold")], border=ERISGOLD, size=11)
tf = textbox(s, Inches(8.05), Inches(3.35), Inches(4.9), Inches(0.65))
para(tf, [("same situation — different verdict: the wording is an attack surface", None)],
     size=10, align=PP_ALIGN.CENTER, color=GRAY, first=True)
block(s, Inches(0.55), Inches(4.35), Inches(12.2), Inches(2.35), "The measurement gap", [
    {"segments": [("Prior work probes one bias at a time and reports a ", None),
                  ("single scalar robustness score", "hi"),
                  (". Across n independent failure directions, a scalar discards n−1 of them. ", None),
                  ("No post-hoc procedure recovers the lost structure.", "gold")]},
], body_size=16)

# ---- 3. our approach --------------------------------------------------------
s = new_slide("Our approach: a geometric vulnerability profile", 3)
tf = textbox(s, Inches(0.55), Inches(1.05), Inches(7.9), Inches(5.6))
para(tf, [("Map each judgment to a point in a ", None), ("7-D harm space", "hi", )],
     size=17, bullet=True, first=True, space_after=2)
para(tf, [("(physical, emotional, financial, autonomy, trust, social, identity; each 0–10).", None)],
     size=13, color=GRAY, space_after=10)
para(tf, [("Apply qualitatively distinct ", None), ("perturbations", "hi"),
          ("; measure ", None), ("displacement", "hi"), (" of the judgment vector.", None)],
     size=17, bullet=True, space_after=10)
para(tf, [("Output a ", None), ("per-model vulnerability profile", "hi"), (", not one number.", None)],
     size=17, bullet=True, space_after=10)
para(tf, [("Runs under a realistic ", None), ("$50/day", "hi"), (" compute budget.", None)],
     size=17, bullet=True)
s.shapes.add_picture(ASSETS + r"\erisml_apple_light.png", Inches(9.7), Inches(1.5), height=Inches(2.4))
tf = textbox(s, Inches(8.9), Inches(4.1), Inches(4.0), Inches(1.6))
para(tf, [("The 7-D space is a measurement-reliable projection of the ", None),
          ("DEME v3", "hi"), (" 9-D moral vector (correspondence table in backup).", None)],
     size=12, align=PP_ALIGN.CENTER, color=GRAY, first=True)

# ---- 4. finding 1 -----------------------------------------------------------
s = new_slide("Finding 1 — Vulnerabilities are selective", 4)
tf = textbox(s, Inches(0.55), Inches(1.0), Inches(5.9), Inches(4.6))
para(tf, [("Displace judgments (real attack surface):", "bold")], size=16, first=True, space_after=6)
para(tf, [("Linguistic framing", "gold")], size=16, bullet=True, space_after=4)
para(tf, [("Emotional anchoring", "gold"), ("  (d\u1d67 = 0.60–1.06)", None)], size=16, bullet=True, space_after=4)
para(tf, [("Irrelevant sensory detail", "gold")], size=16, bullet=True, space_after=12)
para(tf, [("Do ", "bold"), ("not", "em"), (" displace:", "bold")], size=16, space_after=6)
para(tf, [("Gender swap", None)], size=16, bullet=True, space_after=4)
para(tf, [("Evaluation order", None)], size=16, bullet=True)
block(s, Inches(6.8), Inches(1.1), Inches(6.0), Inches(3.3), "The common thread", [
    {"segments": [("The three live surfaces all make morally irrelevant features ", None),
                  ("perceptually salient", "hi"), (".", None)], "space_after": 8},
    {"segments": [("⇒ The attack mechanism is ", None), ("salience manipulation", "gold"),
                  (" — and the framework discriminates it from noise.", None)]},
], body_size=15)
footnote(s, [("We lead with ", None), ("effect sizes", "hi"),
             (" (paired d\u1d67, displacement in harm-points) over σ; empirical control arms rule out stochastic drift.", None)], size=12)

# ---- 5. finding 2 -----------------------------------------------------------
s = new_slide("Finding 2 — Robustness profiles are dissociable", 5)
tf = textbox(s, Inches(0.55), Inches(1.05), Inches(12.2), Inches(2.6))
para(tf, [("No model dominates", "hi"), (" all attack surfaces.", None)],
     size=17, bullet=True, first=True, space_after=10)
para(tf, [("Claude: ", None), ("zero observed sycophancy", "gold"), (" (0/9 — small n, directionally clear) — but ", None), ("worst", "gold"),
          (" emotional-anchoring recovery (20%) and worst divided attention.", None)],
     size=17, bullet=True, space_after=10)
para(tf, [("Flash 2.0: ", None), ("best", "gold"), (" anchoring recovery (73%) — but ", None),
          ("worst", "gold"), (" working memory.", None)],
     size=17, bullet=True)
block(s, Inches(0.55), Inches(4.3), Inches(12.2), Inches(2.2), "Consequence for Secure AI", [
    {"segments": [("Averaging partially-independent dimensions yields a number that ", None),
                  ("describes no model accurately.", "gold")], "space_after": 6},
    {"segments": [("Single-test certification provides false assurance.", "gold")]},
], body_size=16)

# ---- 6. validation ------------------------------------------------------------
s = new_slide("Validation (reviewer-driven) — is the harm space measurable?", 6)
tf = textbox(s, Inches(0.55), Inches(0.95), Inches(12.2), Inches(1.7))
para(tf, [("Is the 7-D space reliably measurable, or asserted?", "bold")], size=17, first=True, space_after=6)
para(tf, [("Inter-model agreement over an ", None), ("open, multi-family panel", "hi"),
          (" (6 models, 5 families: Qwen3, Qwen3-small, GLM-5, GPT-OSS, MiniMax-M2, Gemma) via the ", None),
          ("NRP managed LLM API", "hi"),
          ("; 31 scenarios, disjoint from systems under test, fully reproducible.", None)], size=15)
rows = [
    ["", "ICC(2,k)", "Krippendorff α"],
    ["Financial (highest)", "0.983", "0.902"],
    ["Trust", "0.958", "0.783"],
    ["Physical / Autonomy (lowest)", "0.893", "≈0.57"],
    [[("Overall", "bold")], [("0.969", "gold")], [("0.836", "gold")]],
]
table(s, Inches(3.42), Inches(3.05), Inches(6.5), rows, [3.3, 1.6, 1.6], size=14, row_h=0.36)
footnote(s, [("Test–retest median r = 0.96. Dimensions index a shared structure across model families.", None)],
         y=Inches(5.35), size=12)

# ---- 7. worked exploit ---------------------------------------------------------
s = new_slide("Worked exploit: threshold evasion", 7)
tf = textbox(s, Inches(0.55), Inches(0.95), Inches(12.2), Inches(3.1))
para(tf, [("Content-moderation filter: flag if total harm > T.", "bold")], size=16, first=True, space_after=8)
para(tf, [("Euphemistic rewriting (morally invariant) lowers panel-averaged harm by ", None),
          ("−14.0 points", "gold"),
          (" on the 0–70 scale (4/6 gold > 10 pts); dramatic only +7.3.", None)],
     size=15, bullet=True, space_after=8)
para(tf, [("At a median-calibrated threshold, ", None), ("3 of 6", "gold"),
          (" gold items silently reclassify FLAG → PASS ", None),
          ("(existence proof: n = 6 hand-audited gold items).", "em")],
     size=15, bullet=True, space_after=8)
para(tf, [("Asymmetry", "gold"), (": far easier to ", None), ("evade", "em"),
          (" (hide harm) than to ", None), ("trigger", "em"),
          (" a false positive — the dangerous direction.", None)],
     size=15, bullet=True)
# threshold strip: base (blue) dots pulled left across T by euphemism (orange)
AX_X0, AX_W = 0.8, 11.7


def _ax(v):  # data 0..11.4 -> slide inches
    return Inches(AX_X0 + v * (AX_W / 11.4))


bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(AX_X0), Inches(6.05), Inches(AX_W), Pt(1.4))
bar.fill.solid(); bar.fill.fore_color.rgb = GRAY; bar.line.fill.background()
thr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, _ax(5.6), Inches(4.45), Pt(1.6), Inches(1.75))
thr.fill.solid(); thr.fill.fore_color.rgb = INK; thr.line.fill.background()
tf = textbox(s, _ax(5.6) - Inches(0.8), Inches(4.12), Inches(1.7), Inches(0.3))
para(tf, [("threshold T", "bold")], size=10, align=PP_ALIGN.CENTER, first=True)
tf = textbox(s, _ax(2.2), Inches(4.45), Inches(1.5), Inches(0.3))
para(tf, [("PASS side", None)], size=9, color=GRAY, first=True)
tf = textbox(s, _ax(8.0), Inches(4.45), Inches(1.5), Inches(0.3))
para(tf, [("FLAG side", None)], size=9, color=GRAY, first=True)
ROWS = [(6.3, 4.1, 4.95), (7.2, 4.9, 5.25), (8.4, 5.2, 5.55),
        (7.9, 6.1, 4.95), (9.3, 7.4, 5.25), (10.4, 8.6, 5.55)]
for xb, xe, yy in ROWS:
    larrow(s, _ax(xe), Inches(yy - 0.05), _ax(xb) - _ax(xe), ERISGOLD)
    dot(s, _ax(xb), Inches(yy), ATLASBLUE)
    dot(s, _ax(xe), Inches(yy), ERISGOLD)
tf = textbox(s, Inches(0.85), Inches(4.75), Inches(2.2), Inches(0.6))
para(tf, [("●", "hi"), (" original", None)], size=9, color=GRAY, first=True, space_after=1)
para(tf, [("●", "gold"), (" euphemistic", None)], size=9, color=GRAY)
tf = textbox(s, Inches(6.6), Inches(6.25), Inches(6.0), Inches(0.3))
para(tf, [("schematic geometry — the measured statistics are the bullets above", None)],
     size=9, align=PP_ALIGN.RIGHT, color=GRAY, first=True)

# ---- 8. kernel ---------------------------------------------------------------
s = new_slide("The exploit survives a principled kernel", 8)
tf = textbox(s, Inches(0.55), Inches(0.92), Inches(12.2), Inches(0.45))
para(tf, [("Route each scenario through the real ", None), ("ErisML DEME v3", "hi"), (" decision kernel:", None)],
     size=15, first=True)
BLUE6 = RGBColor(0xEE, 0xF4, 0xFA)
LIGHT = RGBColor(0xF3, 0xF3, 0xF3)
fbox(s, Inches(1.0), Inches(1.95), Inches(1.9), Inches(0.62),
     [("scenario text", None)], fill=BLUE6, size=11)
fbox(s, Inches(3.5), Inches(1.95), Inches(2.3), Inches(0.62),
     [("LLM extracts\n", None), ("ethical facts", "bold")], fill=BLUE6, size=11)
fbox(s, Inches(6.45), Inches(1.95), Inches(2.5), Inches(0.62),
     [("GenevaEMV3\n", "mono"), ("fixed rule module", None)], border=GRAY, fill=LIGHT, size=10)
fbox(s, Inches(9.6), Inches(1.95), Inches(2.9), Inches(0.62),
     [("typed verdict\n", None), ("forbid ⋯ strongly_prefer", "mono")], fill=BLUE6, size=10)
harrow(s, Inches(2.98), Inches(2.18), Inches(0.44))
harrow(s, Inches(5.88), Inches(2.18), Inches(0.44))
harrow(s, Inches(9.03), Inches(2.18), Inches(0.44))
varrow(s, Inches(4.52), Inches(1.36), Inches(0.5), color=ERISGOLD)
tf = textbox(s, Inches(4.95), Inches(1.30), Inches(3.6), Inches(0.5))
para(tf, [("euphemistic rewrite lands ", None), ("HERE", "gold")], size=11, first=True)
tf = textbox(s, Inches(6.45), Inches(2.60), Inches(2.5), Inches(0.3))
para(tf, [("rules unchanged", None)], size=9, align=PP_ALIGN.CENTER, color=GRAY, first=True)
block(s, Inches(0.55), Inches(3.15), Inches(12.2), Inches(2.3), "Across 161 scenario–model cases", [
    {"segments": [("Euphemistic rewriting flips forbid/avoid → permissive in ", None),
                  ("13.7%", "gold"), (" of cases; mean shift ", None), ("+0.65", "gold"),
                  (" ordinal steps.", None)], "bullet": True, "space_after": 8},
    {"segments": [("forbid verdicts fall ", None), ("44 → 27", "gold"),
                  (" (−39%); strongly_prefer rises 46 → 78.", None)], "bullet": True},
], body_size=16)
tf = textbox(s, Inches(0.55), Inches(5.62), Inches(12.2), Inches(1.05))
para(tf, [("A rule engine inherits, rather than cures, the front-end's salience vulnerability.", "gold"),
          ("  Harden ", None), ("fact extraction", "em"), (", not just the rule.", None)],
     size=16, first=True, space_after=4)
para(tf, [("DEME v3 is our own kernel, frozen for this test (nothing tuned to these flips) — a stress substrate, not the ground truth (which is the audited meaning-invariance of the rewrites).", None)],
     size=11, color=GRAY)

# ---- 9. defense (since submission) ---------------------------------------------
s = new_slide("Since submission — a measured defense", 9)
tf = textbox(s, Inches(0.55), Inches(0.90), Inches(12.2), Inches(0.4))
para(tf, [("Second instantiation from here on: validated per-axis learned encoders (", None),
          ("xbse", "mono"), (" — one small trained scorer per axis) replace the LLM judges; same DEME kernel.", None)],
     size=12, color=GRAY, first=True)
# mechanism flow: x -> generated class -> encoders -> average -> decide
BLUE6 = RGBColor(0xEE, 0xF4, 0xFA)
ORANGE6 = RGBColor(0xFC, 0xF1, 0xE8)
TEAL8 = RGBColor(0xE9, 0xF5, 0xF1)
fbox(s, Inches(0.7), Inches(1.55), Inches(1.35), Inches(0.6),
     [("input x", None)], fill=BLUE6, size=11)
fbox(s, Inches(2.6), Inches(1.55), Inches(2.55), Inches(0.6),
     [("paraphrase class\n", None), ("x₁ ⋯ xₘ (generated)", None)],
     border=ERISGOLD, fill=ORANGE6, size=10, dashed=True)
fbox(s, Inches(5.7), Inches(1.55), Inches(1.95), Inches(0.6),
     [("validated encoders", None)], fill=BLUE6, size=11)
fbox(s, Inches(8.2), Inches(1.55), Inches(2.1), Inches(0.6),
     [("average\n", "bold"), ("over the class", None)], border=CBTEAL, fill=TEAL8, size=10)
fbox(s, Inches(10.85), Inches(1.55), Inches(1.4), Inches(0.6),
     [("decide", None)], fill=BLUE6, size=11)
harrow(s, Inches(2.10), Inches(1.77), Inches(0.44))
harrow(s, Inches(5.20), Inches(1.77), Inches(0.44))
harrow(s, Inches(7.70), Inches(1.77), Inches(0.44))
harrow(s, Inches(10.35), Inches(1.77), Inches(0.44))
tf = textbox(s, Inches(2.15), Inches(2.22), Inches(3.6), Inches(0.6))
para(tf, [("the measured hole: a euphemistic x yields a class that ", None),
          ("stays", "em"), (" euphemistic", None)],
     size=9, align=PP_ALIGN.CENTER, color=ERISGOLD, first=True)
tf = textbox(s, Inches(10.15), Inches(2.22), Inches(2.8), Inches(0.6))
para(tf, [("refusal → singleton → ", None), ("escalate", "bold")],
     size=9, align=PP_ALIGN.CENTER, color=GRAY, first=True)
tf = textbox(s, Inches(0.55), Inches(2.95), Inches(12.2), Inches(3.0))
para(tf, [("At scale (60 held-out items, m = 6): raw drift ", None),
          ("0.407 → θ\u2094 = 0.219", "gold"), (" — the mechanism ", None), ("halves", "hi"),
          (" salience drift (θ\u2094: mean per-dimension decision movement). The registered bar θ\u2094 ≤ 0.5 ", None),
          ("binds for adversarial-register inputs", "hi"),
          (" (raw 0.67–0.85); the defended run there is the registered open item.", None)],
     size=15, bullet=True, first=True, space_after=10)
para(tf, [("LLM paraphrasers ", None), ("refuse 24%", "gold"),
          (" of harmful inputs. A non-refusing red-team paraphraser (NLLB back-translation, 6 pivots): ", None),
          ("θ\u2094 = 0.301", "gold"), (" on harmful content — the weakness is the ", None),
          ("generator's", "hi"), (", not the mechanism's.", None)],
     size=15, bullet=True)
footnote(s, [("Natural (not adversarial) paraphrases at scale. Gold-set decision translation (measured negative, n=6): adversarial-register rewrites still flip ", None),
             ("3/3", "gold"),
             (" flagged items defended, displacement −17% — paraphrases of euphemism stay euphemistic; the register gap is the live surface, and escalation carries it. July 2026.", None)], size=10)

# ---- 10. bifactor (since submission) — native editable chart -------------------
s = new_slide("Since submission — reliable ≠ distinct: the geometry is bifactor", 10)
tf = textbox(s, Inches(0.55), Inches(0.92), Inches(12.2), Inches(1.1))
para(tf, [("ICC 0.969 = reliably measurable. Are the dimensions distinct?", "bold"),
          ("  A registered test on the ", None),
          ("fuller DEME 9-axis set + identity_attack", "hi"),
          (" (the moral-foundations superset the 7-D harm space of slides 2–3 projects from — backup B1): a ", None),
          ("general moral-valence channel G", "hi"), (" gate-passes (AUROC ", None),
          ("0.856", "gold"), ("); five named axes are ≥ ", None), ("0.98", "gold"),
          (" predictable from G alone.", None)],
     size=14, first=True)
# native clustered-bar chart: categories bottom-up so purity lands on top
cats = ["autonomy", "environmental", "privacy", "identity attack", "physical",
        "epistemic", "fairness", "care", "loyalty", "legitimacy", "purity"]
cd = CategoryChartData()
cd.categories = cats
cd.add_series("mostly carried by G",
              (None, None, None, None, None, 0.9396, 0.9836, 0.9873, 0.9900, 0.9986, 0.9997))
cd.add_series("mixed", (None, None, None, 0.6658, 0.8197, None, None, None, None, None, None))
cd.add_series("specific", (0.4945, 0.5538, 0.6421, None, None, None, None, None, None, None, None))
gframe = s.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.4), Inches(2.05), Inches(6.4), Inches(4.35), cd)
chart = gframe.chart
chart.has_title = False
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(9)
for ser, col in zip(chart.series, (ATLASBLUE, CBTEAL, ERISGOLD)):
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = col
va = chart.value_axis
va.minimum_scale = 0.4
va.maximum_scale = 1.0
va.tick_labels.font.size = Pt(9)
chart.category_axis.tick_labels.font.size = Pt(9)
# right column: consequence + mapping
block(s, Inches(7.1), Inches(2.05), Inches(5.75), Inches(2.6), "Why this doesn't undercut Findings 1–2", [
    {"segments": [("Harm ", None), ("scores", "em"), (" collapse to G; ", None),
                  ("vulnerability does not", "hi"), (".", None)], "space_after": 6},
    {"segments": [("The attack surfaces that dissociate (Finding 2) are the specific axes that survive G — ", None),
                  ("“G plus physical” discards where manipulation lives", "gold"), (".", None)]},
], body_size=14)
tf = textbox(s, Inches(7.1), Inches(4.85), Inches(5.75), Inches(1.6))
para(tf, [("The ", None), ("judge panel's own 31×7 matrix replicates it", "hi"),
          (": PC1 = ", None), ("54%", "gold"),
          (" of variance, sole factor surviving parallel analysis; per-judge FAs agree (all six models one factor, Tucker congruence 0.989–0.999).", None)],
     size=12, first=True)
footnote(s, [("Bifactor: learned-encoder instantiation (", None), ("xbse", "mono"),
             ("), 12×12 gate at margin 0.05; identity_attack divergence recorded, not adjudicated. Panel FA: Horn's parallel analysis, n=31; per-judge FAs agree — one factor in all six models (PC1 0.51–0.60), Tucker congruence 0.989–0.999 vs consensus. July 2026.", None)], size=10)

# ---- 11. implications ------------------------------------------------------------
s = new_slide("Implications for Secure AI deployment", 11)
tf = textbox(s, Inches(0.55), Inches(1.1), Inches(12.2), Inches(5.2))
para(tf, [("Where to look", "hi"), (": attacks that repackage the same facts with different ", None),
          ("salience", "gold"), (" (euphemistic minimization).", None)],
     size=17, bullet=True, first=True, space_after=12)
para(tf, [("Prompt-level defenses are bounded", "hi"),
          (": explicit warnings recover only ~38% — a ceiling co-occurring with universal overconfidence (ECE, calibration error, 0.19–0.42).", None)],
     size=17, bullet=True, space_after=12)
para(tf, [("Ask the right question", "hi"), (": “which vulnerabilities matter for our use case?” — not “what's the robustness score?”", None)],
     size=17, bullet=True, space_after=12)
para(tf, [("Hardening target", "hi"), (": the fact-extraction front-end, since downstream kernels inherit its failures.", None)],
     size=17, bullet=True)

# ---- 12. limitations ------------------------------------------------------------
s = new_slide("What we have not shown — limitations", 12)
tf = textbox(s, Inches(0.55), Inches(1.05), Inches(12.2), Inches(5.4))
para(tf, [("Adversarial defense is unproven at scale", "hi"),
          (". Class averaging halves natural-paraphrase drift; against adversarial register the gold set (n=6) shows it ", None),
          ("does not hold", "gold"),
          (" (3/3 still flip). Escalate-by-default carries these until a register-crossing generator exists; the at-scale adversarial run is registered, unrun.", None)],
     size=15, bullet=True, first=True, space_after=10)
para(tf, [("Small samples on two claims", "hi"),
          (": the end-to-end exploit is an n=6 existence proof; “zero sycophancy” is 0/9. Directionally clear, not population rates.", None)],
     size=15, bullet=True, space_after=10)
para(tf, [("Two instantiations", "hi"),
          (": attack results are on LLM judges; defense/bifactor on the learned encoders. The native re-run of all five tracks through one pipeline is designed, not yet run.", None)],
     size=15, bullet=True, space_after=10)
para(tf, [("Corpus scope", "hi"),
          (": moderation/social-norm corpora, predominantly English; the contraction's 0.863 is within-corpus-family. Cross-cultural and cross-family transfer untested.", None)],
     size=15, bullet=True)
footnote(s, [("Every number in the talk is hedged at its point of use; this slide collects the four that bound the claims.", None)], size=11)

# ---- 13. philosophy engineering -------------------------------------------------
s = new_slide("The larger program: philosophy engineering", 13)
tf = textbox(s, Inches(0.55), Inches(0.95), Inches(12.2), Inches(1.5))
para(tf, [("This paper is one instance of a method: take a construct usually settled by ", None),
          ("argument", "hi"),
          (" — what is a harm? has this judgment been manipulated? — and build it into an ", None),
          ("instrument", "gold"), (" you can measure, falsify, and harden.", None)],
     size=16, first=True)
tf = textbox(s, Inches(0.55), Inches(2.5), Inches(6.65), Inches(3.9))
para(tf, [("What makes it engineering, not rhetoric", "bold")], size=15, first=True, space_after=10)
para(tf, [("A ", None), ("space", "hi"), (", not a slogan: judgments are points, manipulation is ", None),
          ("displacement", "hi"), (", robustness is ", None), ("invariance", "hi"), (".", None)],
     size=14, bullet=True, space_after=12)
para(tf, [("Pre-registration", "hi"), (" + ", None), ("admission gates that retract", "hi"),
          (": this talk's own bifactor test ", None), ("failed", "em"),
          (" its P1; a channel rescue was ", None), ("refuted", "em"),
          (". Claims that do not transfer are ", None), ("publicly demoted", "gold"), (".", None)],
     size=14, bullet=True)
block(s, Inches(7.45), Inches(2.5), Inches(5.35), Inches(3.3), "The conjecture behind the program (not tested here)", [
    {"segments": [("Normative and aesthetic structure is ", None), ("geometric", "gold"),
                  (" — it lives in ", None), ("compressibility", "hi"), (" and ", None),
                  ("coarse-graining under an observer", "hi"), (".", None)], "space_after": 12},
    {"segments": [("One instrument, many domains: ", None), ("harm", "hi"),
                  (" today; ", None), ("aesthetics, law, cognition", "hi"),
                  (" on the same frame.", None)]},
], body_size=14)
footnote(s, [("A philosophical claim you can ", None), ("build, measure, and be proven wrong about", "gold"),
             (" outranks one you can only defend.", None)])

# ---- 14. takeaways --------------------------------------------------------------
s = new_slide("Takeaways", 14)
block(s, Inches(0.55), Inches(1.15), Inches(12.2), Inches(3.9), "Four takeaways", [
    {"segments": [("1.  Vulnerabilities are ", None), ("selective", "gold"),
                  (" — salience manipulation is the surface.", None)], "space_after": 10},
    {"segments": [("2.  Robustness profiles are ", None), ("dissociable", "gold"),
                  (" — and a single score is now ", None), ("measured", "em"),
                  (" to be mostly general valence.", None)], "space_after": 10},
    {"segments": [("3.  The exploit is ", None), ("end-to-end", "gold"),
                  (" — it flips scalar scores ", None), ("and", "em"), (" typed verdicts.", None)], "space_after": 10},
    {"segments": [("4.  It is ", None), ("defensible", "gold"),
                  (" — class averaging halves natural-paraphrase drift; adversarial-at-scale is the registered open item.", None)]},
], body_size=17)
footnote(s, [("Reproducible: ", None), ("pip install agi-hpc", "mono"), (" (tag ", None),
             ("bds2026-v2", "mono"), (") · ", None), ("moral-spectrum-analyzer", "mono"), (" · ", None),
             ("xbse", "mono"),
             (" · kernel: ErisML / DEME v3 · five tracks, ~8,000 calls under $50/day (Kaggle); validation panel on the open NRP managed LLM API.", None)],
         y=Inches(5.6), size=11)

# ---- 14. thanks --------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
s.shapes.add_picture(ASSETS + r"\atlas_mark_light.png", Inches(5.75), Inches(0.9), height=Inches(1.85))
tf = textbox(s, Inches(1.2), Inches(3.1), Inches(10.93), Inches(3.4))
para(tf, [("Thank you", "hi")], size=30, align=PP_ALIGN.CENTER, first=True, space_after=10)
para(tf, [("Andrew H. Bond    ·    ", None), ("andrew.bond@sjsu.edu", "mono")],
     size=16, align=PP_ALIGN.CENTER, color=INK, space_after=6)
para(tf, [("Selective Invariance Violations in LLM Moral Judgment", None)],
     size=13, align=PP_ALIGN.CENTER, color=GRAY, space_after=2)
para(tf, [("IEEE BigDataService 2026 — Secure AI", None)],
     size=13, align=PP_ALIGN.CENTER, color=GRAY, space_after=12)
para(tf, [("Questions?", "em")], size=13, align=PP_ALIGN.CENTER, color=INK)

# ---- 15. BACKUP: grounding table -------------------------------------------------
s = new_slide("Backup — Grounding: 7-D harm space → DEME v3 9-D vector", "B1")
rows = [
    ["Harm dim.", "DEME v3 axis", "Mapping"],
    ["physical", "physical harm", "direct"],
    ["autonomy", "autonomy / consent", "direct"],
    ["trust", "legitimacy / trust", "direct"],
    ["social", "societal / environmental", "direct"],
    ["financial", "fairness / equity", "partial"],
    ["emotional", "virtue / care", "partial"],
    ["identity", "privacy / standing", "partial"],
    ["—", "rights, epistemic quality", "not scored"],
]
table(s, Inches(2.87), Inches(1.05), Inches(7.6), rows, [2.6, 3.4, 1.6], size=14, row_h=0.34)
tf = textbox(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(1.0))
para(tf, [("The evaluation pipeline is the ", None), ("measurement front-end", "hi"),
          ("; the ErisML compiler's ", None), ("DEME bridge", "hi"),
          (" is the downstream decision kernel (", None), ("DEMEVerdict", "mono"), (").", None)],
     size=14, align=PP_ALIGN.CENTER, first=True)

# ---- 16. BACKUP: discovery -------------------------------------------------------
s = new_slide("Backup — the harm space is extensible and falsifiable", "B2")
tf = textbox(s, Inches(0.55), Inches(0.95), Inches(12.2), Inches(5.3))
para(tf, [("Discovery loop", "bold"),
          (": residual analysis flags candidate missing dimensions → pre-registered admission gate.", None)],
     size=16, first=True, space_after=12)
para(tf, [("Scorecard: 3 flagged → ", None), ("1 validated", "gold"), (" (identity_attack), ", None),
          ("1 retracted", "gold"), (" (threat), ", None), ("1 declined", "gold"),
          (" (sexual content → policy channel); a registered rescue of the rights channel ", None),
          ("refuted", "gold"), (" (trained 0.509 vs null 0.512). ", None),
          ("The gate really rejects — in both directions.", "hi")],
     size=16, bullet=True, space_after=10)
para(tf, [("identity_attack", "mono"), (": cross-dataset held-out AUROC ", None), ("0.80", "gold"),
          (" CI [0.78, 0.83] (+0.25 over null, 2 corpora, n = 6400) — wired live as a ", None),
          ("10th channel", "hi"), (".", None)],
     size=16, bullet=True, space_after=10)
para(tf, [("Cross-lingual at scale: invariance index ", None), ("0.72–0.80", "gold"),
          (" across es/ar/zh/hi/sw, ", None), ("harmful ≈ benign", "hi"), (".", None)],
     size=16, bullet=True)
footnote(s, [("Validated on the learned-encoder perception layer (", None), ("xbse", "mono"),
             ("); native DEME re-run of the five tracks is the extended version. Post-camera-ready results, July 2026.", None)], size=11)

# ---- speaker notes (one entry per slide, same order as above) -------------------
NOTES = [
# 1 — title
"""TIMING: 16 slides, but the live talk is 14 (two backups after Thanks). Plan ~18 min + Q&A: problem 1:30, approach 1:45, findings 3:30, validation 1:15, exploit 2:00, kernel 1:45, defense 1:45, bifactor 1:30, implications 1:15, philosophy 1:00, takeaways 1:00.

This paper was ACCEPTED at IEEE BigDataService 2026, Special Track on Secure AI (both reviews: weak accept).

If you are not Andrew: "I'm presenting on behalf of Andrew Bond, San Jose State University."

ONE-SENTENCE SUMMARY: LLMs that make moral/safety decisions can be manipulated just by REWORDING the input — we built a framework that measures exactly WHICH rewordings move WHICH models, showed the attack works end-to-end against a real moderation pipeline, and (since submission) built and measured a defense and stress-tested the coordinate system itself.

"Geometric" just means: a judgment is a POINT in a multi-dimensional harm space; perturbations MOVE the point (displacement). No deep math on stage.""",

# 2 — problem
"""KEY TERM — "behavioral manipulation vulnerability": the attacker never changes the FACTS, only the PRESENTATION (wording, tone, irrelevant detail). A secure evaluator should give the same verdict either way. These models don't.

Why it matters: LLMs already gate content moderation, triage, legal-intake style decisions. If wording moves the verdict, wording is an attack surface.

BOTTOM BLOCK: prior work tests one bias at a time and reports ONE robustness number. Analogy that lands: a single credit score vs. the full credit report — with n independent failure directions, one scalar discards n−1 of them, unrecoverably.

Don't over-dwell; the punchline slides are 7-10.""",

# 3 — approach
"""HOW IT WORKS (3 steps): (1) every judgment is scored as a 7-dimensional harm vector — physical, emotional, financial, autonomy, trust, social, identity — each 0-10, total 0-70. (2) apply perturbations that should NOT change the moral evaluation. (3) measure how far the vector moved = displacement. Output: a per-model VULNERABILITY PROFILE, not one number.

$50/day: the whole 5-model, ~8,000-call evaluation fits a Kaggle budget — say it out loud, reproducibility is a selling point.

RIGHT SIDE: the 7-D space is the measurement-reliable projection of the DEME v3 9-axis moral vector (ErisML compiler). Two axes (rights, epistemic) were dropped for reliability, disclosed. Full correspondence table is in BACKUP B1 — offer it in Q&A rather than walking it live.

Gloss the word SALIENCE when it first comes up on the next slide: "salience just means what stands out to the model."
""",

# 4 — finding 1
"""THE HEADLINE FINDING of the accepted paper. Vulnerabilities are SELECTIVE, not uniform:

MOVE judgments (all 5 models): linguistic framing (euphemistic vs dramatic rewording), emotional anchoring (paired d_z 0.60-1.06 = medium-to-large), irrelevant sensory detail.
DO NOT move judgments: gender swap, evaluation order.

Common thread (right block): the live surfaces all make morally irrelevant features PERCEPTUALLY SALIENT. The attack mechanism is salience manipulation — and the selectivity proves the framework discriminates real vulnerabilities from noise.

FOOTNOTE (methods honesty): effect sizes are paired d_z (t/sqrt(n)) — say "paired standardized differences" if a statistician asks; every test is compared against replication control arms — apparent 6-sigma violations in early analyses VANISHED under those controls.""",

# 5 — finding 2
"""SECOND FINDING: robustness profiles are DISSOCIABLE — no model dominates.

Concrete contrast: Claude showed ZERO sycophancy in our runs — 0 of 9 fabricated corrections accepted (small n; Wilson CI reaches 30%, corroborated by a confidence-increase signature and 0% control-arm flips) — but the WORST emotional-anchoring recovery (20%) and worst divided attention. Gemini Flash 2.0 recovers from emotional manipulation BEST (73%) — but has the worst working memory.

Consequence: averaging partially-independent dimensions gives a number that describes NO model accurately. Single-test certification = FALSE ASSURANCE.

Anticipate: "which model should I use?" — depends which attack surfaces matter for YOUR use case; that's why the output is a profile, not a ranking. (Slide 10 sharpens this: the scalar's failure is now measured, not just argued.)""",

# 6 — validation
"""Both reviewers' #1 concern: "is the 7-D harm space validated or asserted?" This is the camera-ready answer.

Design: 6 OPEN models from 5 families (Qwen3 x2, GLM-5, GPT-OSS, MiniMax-M2, Gemma) score the full 7-D vector on 31 scenarios, 3 repetitions, via the NRP managed LLM API — fully open, re-runnable. Panel is DISJOINT from the 5 models under test.

TERMS: ICC(2,k) = two-way random, average-measures intraclass correlation; 0.969 is excellent. Krippendorff's alpha = chance-corrected agreement; 0.836 is strong.

CAVEATS to volunteer if pressed: (a) ICC(2,k) is panel-consensus reliability — single-rater ICC(2,1) is ~0.84, still good; (b) six models sharing web-scale training data establishes cross-family consistency, NOT human ground truth — human-rater validation is future work; (c) physical/autonomy lowest (~0.57 alpha) is a near-floor artifact. Test-retest r=0.96 sets the stochastic floor.""",

# 7 — worked exploit
"""Reviewer ask: "show an end-to-end exploit." Walk it as a story:

You run a moderation filter: flag anything with total harm > T. The attacker controls ONLY wording. Rewrite the same facts euphemistically — corporate-liability language, "involuntary separation from the platform" style. Panel-averaged harm drops 14.0 points on the 0-70 scale. At a median-calibrated threshold, 3 of the 6 audited gold scenarios silently flip FLAG -> PASS. Nothing about the event changed.

Say the existence-proof framing OUT LOUD (it's now on the slide): n=6 hand-audited items demonstrate the mechanism end-to-end; the displacement magnitude is established on the 31-scenario panel.

THE ASYMMETRY IS THE SCARY PART: euphemism moves scores twice as much as dramatization (-14 vs +7.3). The filter is far easier to EVADE than to trigger falsely — it fails in the dangerous direction.

If asked "isn't lower harm from euphemism rational inference?": the gold transforms are hand-audited to preserve the stated facts — same events, differently packaged; slide 8 shows the FACTS get extracted differently, which is perception failure, not updating.""",

# 8 — kernel
"""Natural objection: "so don't use a scalar threshold — use a principled rule engine." This slide kills that hope.

Setup: route each scenario through the real ErisML DEME v3 kernel: an LLM extracts structured ethical FACTS, then GenevaEMV3 — a fixed, deterministic rule module — maps facts to a typed verdict (forbid / avoid / neutral / prefer / strongly_prefer).

Result across 161 scenario-model cases (open panel): euphemistic rewriting flips a restrictive verdict to permissive in 13.7% of cases; forbid verdicts drop 44 -> 27 (-39%).

QUOTABLE: "a rule engine INHERITS, rather than cures, the front-end's salience vulnerability." The rules are fine — the FACTS feeding them were corrupted. Hardening must target fact extraction / perception. Sets up the defense.

If asked "isn't this garbage-in-garbage-out by construction?": the point is the QUANTIFICATION — 13.7% and -39% rebut the deployment assumption that a symbolic layer confers robustness, and tell you where the hardening budget goes.""",

# 9 — defense
"""THIS SLIDE HAS FOUR IDEAS — walk them in order, slowly. The audience gets lost here if the instrument switch and the "halves drift BUT items still flip" tension are not resolved out loud.

STEP 0 — THE TRANSITION (say this first, near-verbatim): "Everything you've seen so far measured LLM judges. From this slide on we switch to a second instrument: small validated encoders — the xbse library — measuring the same harm geometry and feeding the same decision kernel. Same coordinates, different sensor." Pause here. If they miss this switch, nothing after it makes sense. These are post-camera-ready results.

STEP 1 — THE IDEA, plain words: "Any single wording is just one way of describing a situation. So don't judge the wording — judge the situation: generate six rewordings, score them all, average, THEN decide." If faces are blank: it's like asking a question six different ways and averaging the answers. (For the ML crowd: semantic randomized smoothing / SmoothLLM with paraphrases — but say plainly there is NO formal certificate; the guarantee is empirical.)

STEP 2 — WHAT IT BUYS, one number: on natural rewordings of 60 held-out items, decision movement drops 0.407 -> 0.219 — HALVED. Define theta_d in one breath: "theta_d just measures how far the final decision moves when the input is reworded; zero = perfectly invariant." Don't dwell on the bar clause; if asked why "meets the 0.5 bar" isn't the headline: natural rewording was already under the bar, so the honest claim is the halving — the bar only matters for adversarial rewrites.

STEP 3 — WHERE IT FAILS. Both failures belong to the paraphrase GENERATOR, not to the averaging:
(a) Refusal: LLM paraphrasers refuse 24% of harmful inputs — exactly the content that matters. A back-translation red-team that refuses nothing gets theta_d = 0.301 on harmful content, so refusal is a generator property. Deployment rule: refusal -> single-member class -> ESCALATE to a human, always.
(b) Register — the footer, and THE tension to resolve out loud: "You may notice the footer seems to contradict what I just said: on the six hand-crafted euphemistic attacks, averaging trimmed movement only 17%, and every flagged item still slipped past the threshold. No contradiction — different input classes. The halving is on NATURAL rewordings. Against a deliberate euphemism, the paraphraser produces paraphrases OF the euphemism — polite versions of a polite lie. The average stays polite. The attack survives." That line — "polite versions of a polite lie" — is the one they'll remember; use it.

STEP 4 — THE LESSON (closes the slide, hands off to implications): "So the mechanism works exactly as far as the generator's diversity reaches — which puts the generator inside the trust boundary. The roadmap fix is a generator that provably crosses registers; until it exists, escalation carries the load, and the audit trail records every class member." Measuring the failure this precisely IS the contribution.

(Adaptive-attacker and bar-circularity questions are pre-answered in the slide-14 Q&A note.)""",

# 10 — bifactor
"""Post-camera-ready. THE SLIDE'S QUESTION: ICC 0.97 showed the dimensions are reliably MEASURABLE — a skeptic asks whether they are DISTINCT. We ran the test against ourselves, pre-registered.

SCOPE + THE NEW RESULT (right column, bottom): the bifactor readout is encoder-instantiation; we then factor-analyzed the JUDGE PANEL'S OWN 31x7 score matrix directly. It REPLICATES the general factor: PC1 = 54% of variance (consensus; 53% pooled), and it is the ONLY component surviving Horn's parallel analysis. Per-axis: physical fully specific on the panel too (R^2 ~ 0 — striking agreement with the encoder residual), emotional (0.86) and trust (0.78) most G-laden as predicted; financial and identity DIVERGE across instantiations — recorded, not adjudicated. n=31, so per-axis loadings are coarse; the factor COUNT is the robust claim.

ROBUSTNESS (if pressed "is that just aggregation?"): run separately on each judge's own 31x7 matrix, ALL SIX models keep exactly one factor (PC1 51-60%), and each model's PC1 is the SAME factor as the consensus — Tucker congruence 0.989-0.999, above the 0.95 factor-equality convention. The general factor lives inside every judge, not just in the average.

RESULTS: a general moral-valence channel G is real and strong (cross-dataset AUROC 0.856, n=51k). Five named axes — purity, legitimacy, loyalty, care, fairness — are >=0.98 predictable from G on independent corpora: their gate-passing transfer is general valence, not their named dimension. The 12x12 specificity gate confirms it independently: specific = environmental, privacy, identity_attack, autonomy, physical(+.08, marginal); demoted = the rest. The pre-registered P1 FAILED (4/11 diagonal-dominant vs required 8) — we publish the failure.

CHART: blue = mostly G, teal = mixed, orange = specific; dashed reference at 0.98.

CONSEQUENCE: a scalar harm score is now MEASURED, not argued, to be mostly G — and any scalar robustness score inherits that collapse by construction. Takeaway 2 upgraded.

If asked "isn't G-dominance guaranteed by construction (valence-signed training pairs)?": yes — G is a maximally strong competitor by design, which is exactly why SURVIVING it is informative; privacy/autonomy/environmental measurably do.
If asked about identity_attack: the gate says specific (+.25), A2 residualization says mixed (G-share 0.67) — the divergence is recorded, not adjudicated; the gate is the registered criterion.""",

# 11 — implications
"""Deployment guidance — four points, one line each:

1. WHERE TO LOOK: the live attack surface is salience repackaging — especially euphemistic minimization. Watch for measured, liability-flavored language hiding severity.
2. PROMPT DEFENSES ARE BOUNDED: telling the model "you are being manipulated" recovers only ~38% — co-occurring with universal overconfidence (ECE 0.19-0.42). System-prompt guardrails fail two times out of three.
3. ASK THE RIGHT QUESTION: not "what's the robustness score?" but "WHICH vulnerabilities matter for OUR use case?" — profiles, not rankings.
4. HARDENING TARGET: the fact-extraction / perception front-end — downstream kernels inherit its failures (slide 8), and slide 9 showed perception CAN be hardened measurably.

ECE = expected calibration error: gap between stated confidence and actual accuracy.""",

# 12 — limitations
"""RIGOR SLIDE — deliver it confidently, not apologetically; naming the four bounds is the point. ~45 sec.

Four bounds: (1) ADVERSARIAL DEFENSE UNPROVEN AT SCALE — natural drift halves, but the n=6 gold set shows it does NOT hold against adversarial register (3/3 flip); escalate-by-default covers those, at-scale adversarial run registered/unrun. (2) SMALL SAMPLES — n=6 exploit, 0/9 sycophancy: directional, not population rates. (3) TWO INSTANTIATIONS — attack on LLM judges, defense/bifactor on encoders; the native single-pipeline re-run is designed not run. (4) CORPUS SCOPE — English/moderation, contraction is within-corpus-family.

This slide pre-empts most of the Q&A: if a reviewer raises any of the four, the answer is "yes — it's on the limitations slide, and here's the plan." That converts an attack into a point of agreement. Do not add new hedges live; these four are the complete set.""",

# 13 — philosophy
"""ONE-MINUTE SLIDE — the zoom-out. Deliver the headline sentence, one beat on each column, move on.

LEFT (all measured, grounded in this talk): the method treats constructs usually settled by argument as INSTRUMENTS — a space (points/displacement/invariance), pre-registration, and admission gates that retract. Evidence from this very talk: the bifactor test FAILED its own P1; the rights-channel rescue was REFUTED. Claims that don't transfer get publicly demoted.

RIGHT (explicitly labeled a conjecture, NOT tested here): the wager that normative and aesthetic structure is geometric — compressibility and coarse-graining under an observer — one instrument across harm, aesthetics, law, cognition.

Q&A DISCIPLINE: if someone pushes on the conjecture ("is this science or manifesto?"), the full answer is: "that's the motivating conjecture of a larger program, not a result — this talk contributes the harm instance and its falsification record." Say exactly that and stop. Do NOT elaborate on aesthetics in a Secure AI Q&A.""",

# 14 — takeaways
"""THE FOUR TAKEAWAYS — read slowly; this is the photographed slide:
1. Vulnerabilities are SELECTIVE — salience manipulation is the attack surface.
2. Robustness profiles are DISSOCIABLE — and a single score is now MEASURED to be mostly general valence (bifactor, slide 10).
3. The exploit is END-TO-END — it flips scalar thresholds AND typed rule-engine verdicts.
4. It is DEFENSIBLE — class averaging halves natural-paraphrase drift; adversarial-at-scale is the registered open item.

Attack -> measurement -> defense -> instrument audit. That's the arc.

Reproducibility (footer, one breath): pip install agi-hpc (tag bds2026-v2), moral-spectrum-analyzer, xbse; kernel ErisML/DEME v3; $50/day on Kaggle; validation panel on the open NRP API.

Spoken closer after takeaway 4 (borrowed from slide 12): "One instrument today — many domains tomorrow. Thank you." """,

# 15 — thanks / Q&A prep
"""Q&A PREP — hardest questions, one-line answers:

Q: You showed ICC 0.97, then a bifactor result saying most axes are one factor — measured on different encoders. Did you factor-analyze the panel's own scores? A: Yes — the panel's own 31x7 matrix gives PC1 = 54% of variance, the only component surviving parallel analysis, with physical fully specific there too. The general factor replicates across both instantiations; two per-axis assignments (financial, identity) diverge and are recorded as such.
Q: Isn't G-dominance guaranteed by construction? A: G is a maximally strong competitor by design — which is why surviving it is informative; privacy, autonomy, environmental measurably do.
Q: Your drift bar 0.5 is met by the raw pipeline (0.407). What did the defense buy? A: On natural paraphrases the bar doesn't bind — it was set against adversarial transforms (raw 0.67-0.85); the at-scale result is the HALVING, confirmed on harmful content via a non-refusing generator; defended-adversarial is the registered open item.
Q: With averaging on, how many gold items still flip? What does 13.7% become? A: Measured, and it's an honest negative — on the deployed encoder instrument all 3 flagged gold items STILL flip defended (displacement only −17% vs −46% on natural paraphrases). Paraphrases of euphemism stay euphemistic: the class inherits the register, and register is the attack vector. That's why the residual surface is the generator's register diversity and why escalate-by-default is load-bearing. The DEME-instantiation defended-13.7% is the native re-run.
Q: Euphemism deletes information — isn't a lower score rational? A: Gold transforms are hand-audited to preserve stated facts; the DEME leg shows the same facts get EXTRACTED differently — perception failure, not updating.
Q: Six LLMs agreeing isn't validity — where are humans? A: Agreed — the panel establishes reliability and cross-family consistency; single-rater ICC ~0.84; human-rater validation is future work.
Q: Rule-engine result is garbage-in-garbage-out by construction? A: The quantification is the finding — 13.7% flips and -39% forbid rebut the assumption that a symbolic layer confers robustness, and locate the hardening budget.
Q: identity_attack: gate says specific, residualization says 2/3 G, and it dominates your contraction — is the contraction re-measuring valence? A: The divergence is recorded, not adjudicated; the +0.084 out-of-fold lift is against a baseline already containing the G-heavy axes, which bounds how much can be G-duplication.
Q: Isn't this just prompt sensitivity? A: Prompt sensitivity is undirected noise; this is DIRECTED, selective, reproducible displacement under meaning-preserving transforms, with control arms and a decision consequence.
Q: Model coverage small (5 models / 2 families)? A: Acknowledged; the validation panel adds 6 open models / 5 families; trend claims flagged as suggestive.

Contact: andrew.bond@sjsu.edu""",

# 15 — backup: grounding
"""BACKUP — show if asked "where do the 7 dimensions come from?"

They are a PROJECTION of the DEME v3 9-axis moral vector: four map directly, three partially, two (rights, epistemic quality) are NOT scored — they produced unreliable judgments in preliminary testing, disclosed rather than hidden. Key phrase: "a reliability-driven projection, not a bijection."

If pressed on "identity" mapping weakly to privacy/standing: it is properly fixed by the validated identity_attack channel (backup B2) — now a separately-scored 10th channel in the extended instrument.""",

# 16 — backup: discovery
"""BACKUP — show if asked "is the dimension set fixed/asserted?" or "what's the discovery loop?"

THE LOOP: residual analysis asks "is there moderation signal the current axes don't carry?" -> flags candidates -> each faces a PRE-REGISTERED admission gate (cross-dataset transfer, margin over null, fuzz test).

SCORECARD — the rejections are the credibility: 3 flagged. 1 VALIDATED: identity_attack, held-out AUROC 0.80 CI [0.78-0.83], two corpora, n=6400 — live as a 10th channel. 1 RETRACTED: threat (failed balanced resample). 1 DECLINED: sexual content — a PLATFORM POLICY signal, not a moral axis (consensual explicitness isn't a moral violation; the moral weight was harassment, already covered). PLUS: a registered rescue of the rights channel was REFUTED (trained 0.509 vs null 0.512) — it stays a hand-specified hard channel. The gate really rejects, in both directions.

Also: cross-lingual invariance at scale 0.72-0.80 across es/ar/zh/hi/sw, harmful ≈ benign.

HONESTY LINE: these validations are on the learned-encoder perception layer (xbse); re-running the paper's five tracks natively through this instrument IS the extended version.
If asked "isn't identity_attack just a Perspective API label?": the LABEL exists; the contribution is the INSTRUMENT — discovery from residuals, pre-registered admission, provenance-carrying deployment. Two scope notes we volunteer: the gate leg is cross-dataset structure over a jointly-trained pair (strict train-A/test-B and a third corpus are open), and flag/validation corpora overlap at the family level.""",
]

assert len(NOTES) == len(prs.slides._sldIdLst), (
    f"NOTES has {len(NOTES)} entries but deck has "
    f"{len(prs.slides._sldIdLst)} slides — keep them aligned!")

for _slide, _note in zip(prs.slides, NOTES):
    _slide.notes_slide.notes_text_frame.text = _note

prs.save(OUT)
print(f"saved {OUT} with {len(prs.slides._sldIdLst)} slides and {len(NOTES)} aligned notes")
