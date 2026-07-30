#!/usr/bin/env python3
"""Omron (Kyoto) briefing deck, built from the accepted BDS 2026 paper 0204.

Audience: Omron technical staff + executives. Purpose: credibility and
relationship building, not a specific commercial ask.

Design: the first ten slides stand alone for an executive slot of about
30 minutes. The appendix carries the validation evidence for a longer technical
session, so one file serves both agendas.

Every number here is taken from SecureAI_Bond_2026_slides.tex / the accepted
paper. Nothing is rounded up and nothing is invented.

    python build_omron_pptx.py
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = r"C:\source\agi-hpc\benchmarks\ieee_bds_2026\SecureAI_Omron_Kyoto.pptx"

# Validated ColorBrewer palette, same as the Fukuoka deck ---------------------
BLUE = RGBColor(0x21, 0x71, 0xB5)   # structure / general
ORANGE = RGBColor(0xD9, 0x5F, 0x02)  # attack / specific
TEAL = RGBColor(0x1B, 0x9E, 0x77)    # defended / mitigated
INK = RGBColor(0x1A, 0x23, 0x40)
BLOCKBG = RGBColor(0xE4, 0xED, 0xF7)
GRAY = RGBColor(0x63, 0x63, 0x63)
BLACK = RGBColor(0x20, 0x20, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW = prs.slide_width


def tb(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def runs(p, segs, size, base=BLACK):
    """segs: list of str or (str, color) or (str, color, bold)."""
    for seg in segs:
        if isinstance(seg, str):
            text, color, bold = seg, base, False
        elif len(seg) == 2:
            text, color, bold = seg[0], seg[1], False
        else:
            text, color, bold = seg
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.color.rgb = color
        r.font.bold = bold


def para(tf, segs, size=16, bullet=False, first=False, space=8, indent=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space)
    if indent:
        p.level = indent
    if bullet:
        segs = [("\u2022  ", GRAY)] + list(segs)
    runs(p, segs, size)
    return p


def title_bar(slide, text, num=None):
    tf = tb(slide, Inches(0.62), Inches(0.42), Inches(11.4), Inches(0.7))
    p = tf.paragraphs[0]
    runs(p, [(text, INK, True)], 27)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62),
                                  Inches(1.12), Inches(12.1), Pt(2.2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLUE
    rule.line.fill.background()
    rule.shadow.inherit = False
    if num is not None:
        n = tb(slide, Inches(12.35), Inches(6.92), Inches(0.6), Inches(0.3))
        pp = n.paragraphs[0]
        pp.alignment = PP_ALIGN.RIGHT
        runs(pp, [(str(num), GRAY)], 11)


def block(slide, x, y, w, h, heading, lines, size=15):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = BLOCKBG
    box.line.color.rgb = BLUE
    box.line.width = Pt(1.1)
    box.shadow.inherit = False
    box.adjustments[0] = 0.04
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.11)
    p = tf.paragraphs[0]
    p.space_after = Pt(6)
    runs(p, [(heading, INK, True)], size + 1)
    for ln in lines:
        para(tf, ln, size=size, space=5)
    return box


def footnote(slide, segs, y=Inches(6.82), size=11):
    tf = tb(slide, Inches(0.62), y, Inches(11.6), Inches(0.5))
    runs(tf.paragraphs[0], segs, size, base=GRAY)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def new(title=None, num=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        title_bar(s, title, num)
    return s


def chip(slide, x, y, w, h, text, border, fill=None, size=12, bold=False,
         tcolor=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = WHITE
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(1.1)
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.10
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    runs(p, [(text, tcolor or INK, bold)], size)
    return sh


def arrow(slide, x, y, w, color=GRAY):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.13))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False


def table(slide, x, y, w, rows, widths, size=13, header=True):
    n_r, n_c = len(rows), len(rows[0])
    h = Inches(0.34) * n_r
    shp = slide.shapes.add_table(n_r, n_c, x, y, w, h)
    t = shp.table
    for j, cw in enumerate(widths):
        t.columns[j].width = cw
    for i, row in enumerate(rows):
        for j, cell_val in enumerate(row):
            c = t.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = BLOCKBG if (header and i == 0) else WHITE
            c.margin_left = c.margin_right = Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.02)
            tf = c.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            if isinstance(cell_val, tuple):
                runs(p, [cell_val], size)
            else:
                runs(p, [(cell_val, INK, header and i == 0)], size)
    return t


# ===========================================================================
# 1. Title
# ===========================================================================
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.28), SW, Pt(3))
band.fill.solid()
band.fill.fore_color.rgb = BLUE
band.line.fill.background()
band.shadow.inherit = False

tf = tb(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(1.1))
runs(tf.paragraphs[0],
     [("When AI judges the wording, not the facts", INK, True)], 38)

tf = tb(s, Inches(0.9), Inches(2.62), Inches(11.5), Inches(0.9))
para(tf, [("A measurement instrument for AI judgment consistency, "
           "and a failure mode that survives the obvious fix", GRAY)],
     size=19, first=True)

tf = tb(s, Inches(0.9), Inches(4.05), Inches(11.5), Inches(1.7))
para(tf, [("Prepared for OMRON Corporation", INK, True)], size=17, first=True)
para(tf, [("Kyoto", GRAY)], size=15)
para(tf, [("Andrew H. Bond", INK)], size=15, space=3)
para(tf, [("San Jose State University \u00b7 IEEE Senior Member", GRAY)],
     size=13)

tf = tb(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.6))
para(tf, [("Based on ", GRAY),
          ("IEEE BigDataService 2026", INK, True),
          (", Special Track on Secure AI, paper 0204 (accepted). "
           "All figures measured; sample sizes stated throughout.", GRAY)],
     size=12, first=True)
notes(s, """
Opening frame, roughly 60 seconds. Thank them for the time and say plainly what
this is: a piece of accepted academic work, brought because the failure mode it
found is a safety-engineering problem rather than an academic curiosity.

Say early that the talk includes what did NOT work. That is deliberate and it is
the part most useful to an engineering organisation.

If asked why Omron: because you build instruments that make physical properties
measurable and controllable, and this is the same move applied to a behavioural
property of AI systems.
""")

# ===========================================================================
# 2. The whole story in one page
# ===========================================================================
s = new("The finding, in one page", 2)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(6.35), Inches(5.2))
para(tf, [("1.  A reworded input changes the verdict.", INK, True)],
     size=17, first=True, space=4)
para(tf, [("Rewriting a situation in softer language, with the facts "
           "unchanged, lowers the measured harm score by ", GRAY),
          ("14.0 points on a 0\u201370 scale", ORANGE, True),
          (". Three of six audited items silently move from FLAG to PASS.",
           GRAY)], size=14, space=12)

para(tf, [("2.  It is a systematic fault, not noise.", INK, True)],
     size=17, space=4)
para(tf, [("Control perturbations that should not matter, and do not, "
           "confirm it: swapping gender or reordering cases moves nothing. "
           "Only changes that make irrelevant detail ", GRAY),
          ("more salient", INK, True), (" move the verdict.", GRAY)],
     size=14, space=12)

para(tf, [("3.  The obvious fix does not close it.", INK, True)],
     size=17, space=4)
para(tf, [("Paraphrase-and-average halves the drift on natural rewordings "
           "(0.407 \u2192 0.219). Against deliberately softened register it "
           "fails: ", GRAY),
          ("3 of 3", ORANGE, True),
          (" flagged items still flip. Paraphrases of a euphemism stay "
           "euphemistic.", GRAY)], size=14, space=12)

para(tf, [("4.  So the safe default has to carry it.", INK, True)],
     size=17, space=4)
para(tf, [("Escalate to a human when the check cannot be trusted, rather "
           "than passing. In safety terms, a defined safe state rather than "
           "fail-operational.", GRAY)], size=14)

block(s, Inches(7.25), Inches(1.42), Inches(5.45), Inches(2.5),
      "Why this matters for a components business",
      [[("A text classifier that can be reliably evaded by rewording is a ",
         GRAY), ("systematic", INK, True),
        (" fault with low diagnostic coverage.", GRAY)],
       [("A single robustness score does not detect it, and we measure why "
         "that score is mostly one general dimension.", GRAY)]], size=14)

block(s, Inches(7.25), Inches(4.12), Inches(5.45), Inches(2.5),
      "What we brought",
      [[("An instrument that profiles ", GRAY),
        ("where", INK, True),
        (" a given model is vulnerable, across seven harm axes, validated "
         "for reliability across six models and five families.", GRAY)],
       [("Not a score. A profile you can act on.", INK, True)]], size=14)
notes(s, """
This is the slide to leave on screen if the meeting is cut short. Everything
else expands one of these four points.

Executives usually stop at point 1 and point 4. Engineers go to point 3.

If someone asks "is 14 points a lot?" the answer is that it is the distance that
moved three of six hand-audited items across a median-calibrated decision
threshold. The unit is arbitrary; crossing the threshold is not.
""")

# ===========================================================================
# 3. Where AI now sits in the decision path
# ===========================================================================
s = new("Where AI now sits in the decision path", 3)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(6.4), Inches(2.4))
para(tf, [("Language models increasingly sit ", GRAY),
          ("inside", INK, True),
          (" decisions rather than beside them. They screen content, "
           "triage cases, summarise incidents, and classify reports.",
           GRAY)], size=16, first=True)
para(tf, [("When a model gates a decision, it becomes a component with a "
           "failure mode, and it inherits the obligations of one.", GRAY)],
     size=16)
para(tf, [("A sound evaluator should judge the ", GRAY),
          ("facts", INK, True), (", not their ", GRAY),
          ("presentation", ORANGE, True), (".", GRAY)], size=16)

# flow: same facts -> two wordings -> two verdicts
cx = Inches(7.15)
chip(s, cx + Inches(1.35), Inches(1.62), Inches(2.5), Inches(0.52),
     "the same situation", INK, fill=RGBColor(0xF4, 0xF4, 0xF4), size=13,
     bold=True)
chip(s, cx, Inches(2.72), Inches(2.35), Inches(0.52), "stated plainly", GRAY)
chip(s, cx + Inches(2.75), Inches(2.72), Inches(2.35), Inches(0.52),
     "softened wording", GRAY)
chip(s, cx + Inches(0.5), Inches(3.75), Inches(1.35), Inches(0.5), "FLAG",
     BLUE, size=13, bold=True, tcolor=BLUE)
chip(s, cx + Inches(3.25), Inches(3.75), Inches(1.35), Inches(0.5), "PASS",
     ORANGE, size=13, bold=True, tcolor=ORANGE)
tf = tb(s, cx, Inches(4.42), Inches(5.1), Inches(0.8))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
runs(p, [("same facts, different verdict:\nthe wording is the attack surface",
          GRAY)], 12)

block(s, Inches(0.62), Inches(4.35), Inches(6.4), Inches(2.15),
      "The measurement gap we set out to close",
      [[("Existing work probes one bias at a time and reports a ", GRAY),
        ("single robustness number", INK, True), (".", GRAY)],
       [("Across several independent failure directions, one number keeps one "
         "and discards the rest. No later step recovers them.", GRAY)]],
      size=14)
notes(s, """
Keep this short for executives; it is setup, not content.

The point that lands with a components company: once a model gates a decision it
is a component, and components get characterised, rated, and given failure
modes. Nobody ships a safety relay with one number called "reliability".

If they mention their own deployments here, let them talk. That is more valuable
than the rest of the deck.
""")

# ===========================================================================
# 4. The failure, measured
# ===========================================================================
s = new("The failure, measured", 4)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.8))
para(tf, [("Setup: a content check that flags a case when total measured harm "
           "exceeds a threshold. We rewrite each case in softer language "
           "while holding the facts fixed, then re-measure.", GRAY)],
     size=15, first=True)

rows = [["", "Effect on the measured harm score", "Interpretation"],
        ["Softened wording", ("\u221214.0 points (0\u201370 scale)",
                             ORANGE, True),
         "4 of 6 audited items move more than 10 points"],
        ["Dramatised wording", ("+7.3 points", INK, False),
         "Raising alarm is much harder than lowering it"],
        ["At a calibrated threshold", ("3 of 6 flip FLAG \u2192 PASS",
                                       ORANGE, True),
         "Hand-audited gold items; an existence proof, not a rate"]]
table(s, Inches(0.62), Inches(2.42), Inches(12.1), rows,
      [Inches(2.6), Inches(4.1), Inches(5.4)])

block(s, Inches(0.62), Inches(4.15), Inches(5.9), Inches(2.35),
      "The asymmetry is the dangerous part",
      [[("Hiding harm is far easier than manufacturing it.", INK, True)],
       [("An attacker wants exactly the easy direction, and a reviewer "
         "checking for false alarms will not see it.", GRAY)]], size=14)

block(s, Inches(6.82), Inches(4.15), Inches(5.9), Inches(2.35),
      "Effect sizes, not just significance",
      [[("Emotional anchoring: paired effect size ", GRAY),
        ("d\u2082 = 0.60 to 1.06", INK, True), (".", GRAY)],
       [("Reported as displacement in harm points and effect size, so the "
         "magnitude is legible rather than only the p-value.", GRAY)]],
      size=14)
footnote(s, [("n = 6 hand-audited gold items for the threshold result. "
              "We report it as an existence proof and say so on the limits "
              "slide.", GRAY)])
notes(s, """
This is the evidence slide for the executive half. Spend time here.

Two numbers to say out loud: minus 14 points, and three of six crossing the
threshold. The second is the one that matters, because it is a decision changing,
not a score moving.

The asymmetry line is worth pausing on. It is the reason this is a security
finding and not a quality finding.

Anticipated question: "why only six items?" Answer honestly: those are hand-
audited gold items where we verified the rewrite preserved the facts. Scaling
that audit is exactly the work we would want to do with a partner.
""")

# ===========================================================================
# 5. Systematic, not random
# ===========================================================================
s = new("This is a systematic fault, not noise", 5)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.6))
para(tf, [("The distinction matters because the two call for different "
           "treatment. We tested it directly by including perturbations that "
           "ought to change nothing.", GRAY)], size=15, first=True)

block(s, Inches(0.62), Inches(2.22), Inches(5.9), Inches(2.05),
      "Moves the verdict",
      [[("Linguistic framing", ORANGE, True)],
       [("Emotional anchoring", ORANGE, True)],
       [("Irrelevant sensory detail", ORANGE, True)]], size=15)
block(s, Inches(6.82), Inches(2.22), Inches(5.9), Inches(2.05),
      "Does not move the verdict",
      [[("Swapping the gender of the people involved", GRAY)],
       [("Changing the order cases are evaluated in", GRAY)]], size=15)

block(s, Inches(0.62), Inches(4.5), Inches(12.1), Inches(2.0),
      "What the contrast tells us",
      [[("The three live surfaces share one mechanism: they make morally "
         "irrelevant features ", GRAY), ("more noticeable", INK, True),
        (". The two inert ones do not.", GRAY)],
       [("So the effect is directional and reproducible, which makes it a "
         "systematic fault. It will not average out across many cases, and it "
         "will not be caught by testing more of the same inputs.", GRAY)]],
      size=15)
notes(s, """
This slide is aimed squarely at safety engineers, and it is the one that earns
technical credibility.

The vocabulary to use: systematic versus random. A random fault is diluted by
volume and caught by statistics. A systematic fault is reproducible, so volume
does not help and a monitoring approach tuned for random error has almost no
diagnostic coverage for it.

The control arms are the methodological point. Without them, someone can always
say the effect is stochastic drift in the model. Gender swap and evaluation
order are the arms that rule that out.
""")

# ===========================================================================
# 6. What does not fix it
# ===========================================================================
s = new("What does not fix it, measured honestly", 6)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.55))
para(tf, [("The natural defence: generate several paraphrases of the input, "
           "score them all, and average. If wording is the attack surface, "
           "averaging over wordings should blunt it.", GRAY)],
     size=15, first=True)

y = Inches(2.2)
chip(s, Inches(0.62), y, Inches(1.5), Inches(0.62), "input", BLUE,
     fill=BLOCKBG)
arrow(s, Inches(2.22), y + Inches(0.24), Inches(0.5))
chip(s, Inches(2.82), y, Inches(2.3), Inches(0.62), "paraphrase set",
     ORANGE)
arrow(s, Inches(5.22), y + Inches(0.24), Inches(0.5))
chip(s, Inches(5.82), y, Inches(1.9), Inches(0.62), "score each", BLUE)
arrow(s, Inches(7.82), y + Inches(0.24), Inches(0.5))
chip(s, Inches(8.42), y, Inches(1.7), Inches(0.62), "average", TEAL,
     fill=RGBColor(0xE8, 0xF6, 0xF1))
arrow(s, Inches(10.22), y + Inches(0.24), Inches(0.5))
chip(s, Inches(10.82), y, Inches(1.9), Inches(0.62), "decide", BLUE)
tf = tb(s, Inches(2.82), y + Inches(0.7), Inches(2.3), Inches(0.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
runs(p, [("the hole is here", ORANGE, True)], 11)

block(s, Inches(0.62), Inches(3.62), Inches(5.9), Inches(1.5),
      "On natural rewordings, it works",
      [[("Drift halves: ", GRAY), ("0.407 \u2192 0.219", TEAL, True),
        (", over 60 held-out items with six paraphrases each.", GRAY)]],
      size=14)
block(s, Inches(6.82), Inches(3.62), Inches(5.9), Inches(1.5),
      "On deliberately softened register, it does not",
      [[("3 of 3", ORANGE, True),
        (" flagged items still flip with the defence active. Displacement "
         "falls only 17 percent, against 46 percent for natural rewording.",
         GRAY)]], size=14)

block(s, Inches(0.62), Inches(5.35), Inches(12.1), Inches(1.15),
      "Why, and this is the transferable lesson",
      [[("A paraphrase of a euphemism is still a euphemism. The generator "
         "stays inside the register it was given, so averaging over it "
         "averages over the wrong set. ", GRAY),
        ("Register is the attack vector, and the defence has to cross it "
         "deliberately.", INK, True)]], size=14)
footnote(s, [("Related measurement: the paraphrase generators refuse 24 "
              "percent of harmful inputs outright. A non-refusing "
              "back-translation generator reaches drift 0.301, so that "
              "weakness is the generator's rather than the mechanism's.",
              GRAY)])
notes(s, """
This is the slide that distinguishes this work from a vendor pitch, and it is
the one to be proud of. Present it as a negative result you went looking for.

The line to deliver slowly: a paraphrase of a euphemism is still a euphemism.
Engineers recognise that immediately as a coverage argument about the generator,
not a tuning problem.

If they ask whether the defence is therefore useless: no. It halves natural
drift, which is real. It just cannot be relied on for the adversarial case, so
it must not be the last line.

The at-scale adversarial run is registered and not yet run. Say so if pressed.
""")

# ===========================================================================
# 7. What helps today
# ===========================================================================
s = new("What helps today", 7)
rows = [["Measure", "Effect", "Status"],
        ["Escalate to a human when the check is untrustworthy",
         ("carries the cases the defence misses", TEAL, True),
         "recommended today"],
        ["Average over a paraphrase set",
         ("halves drift on natural rewording", TEAL, False),
         "useful, not sufficient"],
        ["Warn the model in the prompt",
         ("recovers about 38 percent only", ORANGE, False),
         "bounded"],
        ["Harden the fact-extraction front end",
         ("removes the cause rather than the symptom", INK, False),
         "the real target"]]
table(s, Inches(0.62), Inches(1.5), Inches(12.1), rows,
      [Inches(4.6), Inches(4.4), Inches(3.1)])

block(s, Inches(0.62), Inches(3.55), Inches(5.9), Inches(1.55),
      "The safety-engineering read",
      [[("Since the fault is systematic and the mitigation is partial, the "
         "system needs a ", GRAY), ("defined safe state", INK, True),
        (" rather than a best-effort pass.", GRAY)]], size=14)
block(s, Inches(6.82), Inches(3.55), Inches(5.9), Inches(1.55),
      "Prompt-level warnings hit a ceiling",
      [[("They co-occur with overconfidence: calibration error runs 0.19 to "
         "0.42 across models. The model does not know it has been moved.",
         GRAY)]], size=14)

block(s, Inches(0.62), Inches(5.32), Inches(12.1), Inches(1.18),
      "Where the cause actually lives",
      [[("Downstream decision logic inherits whatever the front end got "
         "wrong. If the extraction of facts from text is movable by wording, "
         "every rule built on top of it is movable too. That is where we "
         "would put engineering effort.", GRAY)]], size=14)
notes(s, """
Executives want this slide. It is the "so what do we do" slide.

Order matters: escalate-by-default is available today and costs only human
review capacity. Front-end hardening is the real fix but it is engineering work.

The calibration number is worth one sentence: the model is confidently wrong
after being manipulated, which is why you cannot use its own confidence as the
trigger for escalation.

If they ask for a recommendation in one line: do not let a text-only check be
the last gate on a consequential decision.
""")

# ===========================================================================
# 8. Why one robustness score is unsafe
# ===========================================================================
s = new("Why a single robustness score is unsafe", 8)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.55))
para(tf, [("We compared models across attack surfaces. No model wins "
           "everywhere, and the trade-offs are not visible in any single "
           "number.", GRAY)], size=15, first=True)

rows = [["Model", "Strong", "Weak"],
        ["Claude", ("no sycophancy observed (0 of 9)", TEAL, True),
         ("worst recovery from emotional anchoring, 20 percent", ORANGE,
          False)],
        ["Flash 2.0", ("best anchoring recovery, 73 percent", TEAL, True),
         ("worst working memory", ORANGE, False)]]
table(s, Inches(0.62), Inches(2.15), Inches(12.1), rows,
      [Inches(2.0), Inches(5.05), Inches(5.05)])

block(s, Inches(0.62), Inches(3.5), Inches(5.9), Inches(1.7),
      "The consequence for procurement",
      [[("Averaging partly independent dimensions produces a number that "
         "describes ", GRAY), ("no model accurately", INK, True),
        (". Certifying on one test gives false assurance.", GRAY)]], size=14)
block(s, Inches(6.82), Inches(3.5), Inches(5.9), Inches(1.7),
      "And we measured why",
      [[("A follow-up test found one general dimension explains most of the "
         "spread: five named axes are 98 percent predictable from it. A "
         "scalar score is mostly that one channel.", GRAY)]], size=14)

block(s, Inches(0.62), Inches(5.42), Inches(12.1), Inches(1.08),
      "The question to ask a supplier",
      [[("Not \u201cwhat is your robustness score\u201d but ", GRAY),
        ("\u201cwhich vulnerabilities does your model have, and do they "
         "overlap with the decisions we are giving it\u201d", INK, True),
        (".", GRAY)]], size=15)
notea = """
For executives this is the procurement slide, and it is the most commercially
useful one in the deck.

The Claude and Flash rows are the memorable illustration: the model with the
cleanest record on one axis is the worst on another. A single ranking hides
exactly the information you need.

Caveat honestly: zero sycophancy is 0 of 9, so it is directional, not a rate.
Say that before anyone asks.

The bifactor result is on the appendix slide if they want the method.
"""
notes(s, notea)

# ===========================================================================
# 9. What the instrument measures
# ===========================================================================
s = new("What the instrument measures", 9)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.55))
para(tf, [("Each judgment becomes a point in a seven-axis harm space. A "
           "perturbation moves the point. The distance moved is the "
           "vulnerability, per axis, per model.", GRAY)], size=15,
     first=True)

axes = ["physical", "emotional", "financial", "autonomy", "trust", "social",
        "identity"]
x0 = Inches(0.62)
for i, ax in enumerate(axes):
    chip(s, x0 + Emu(int(Inches(1.72).emu * i)), Inches(2.18),
         Inches(1.6), Inches(0.55), ax, BLUE, fill=BLOCKBG, size=12)
tf = tb(s, Inches(0.62), Inches(2.8), Inches(12.1), Inches(0.4))
runs(tf.paragraphs[0], [("each scored 0 to 10, so the total runs 0 to 70",
                         GRAY)], 12)

block(s, Inches(0.62), Inches(3.35), Inches(3.85), Inches(1.85),
      "Output is a profile",
      [[("Not one number. A per-model map of which surfaces move it and by "
         "how much.", GRAY)]], size=14)
block(s, Inches(4.75), Inches(3.35), Inches(3.85), Inches(1.85),
      "Validated for reliability",
      [[("Agreement across six models in five families: ", GRAY),
        ("ICC 0.969", INK, True), (", test\u2013retest 0.96.", GRAY)]],
      size=14)
block(s, Inches(8.88), Inches(3.35), Inches(3.84), Inches(1.85),
      "Runs on a modest budget",
      [[("About 8,000 model calls, under ", GRAY),
        ("50 US dollars per day", INK, True),
        (". Reproducible on an open API.", GRAY)]], size=14)

block(s, Inches(0.62), Inches(5.42), Inches(12.1), Inches(1.08),
      "The engineering claim",
      [[("A property usually settled by argument, namely whether a judgment "
         "has been manipulated, becomes something you can ", GRAY),
        ("measure, monitor, and set a threshold on", INK, True),
        (". That is the same move as putting a sensor on a process.", GRAY)]],
      size=15)
notes(s, """
This is the slide where the Omron parallel is explicit, and it is worth stating
directly rather than leaving implied: their business is making physical
properties measurable and then controllable. This is that move applied to a
behavioural property of an AI system.

If the room is receptive, this is the natural place to ask where in their
products a text-based judgment currently gates something.

The seven axes are worth reading aloud once. They are recognisable categories,
not jargon, which is part of why the reliability numbers came out as high as
they did.
""")

# ===========================================================================
# 10. Closing: capability + invitation
# ===========================================================================
s = new("Where this could be useful to Omron", 10)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.5))
para(tf, [("We are not proposing a specific engagement today. We are "
           "offering a capability and asking where it touches your roadmap.",
           GRAY)], size=15, first=True)

block(s, Inches(0.62), Inches(2.1), Inches(3.85), Inches(2.5),
      "What we can do now",
      [[("Profile a model or pipeline across the seven axes and report where "
         "it is movable.", GRAY)],
       [("Build a register-crossing test set for a specific domain.", GRAY)],
       [("Advise on where the safe-state boundary belongs.", GRAY)]],
      size=14)
block(s, Inches(4.75), Inches(2.1), Inches(3.85), Inches(2.5),
      "Questions we would ask you",
      [[("Where does a text-based judgment currently gate a decision in your "
         "products or operations?", GRAY)],
       [("Which harm axes matter for those decisions?", GRAY)],
       [("What is the cost of a missed flag versus a false alarm?", GRAY)]],
      size=14)
block(s, Inches(8.88), Inches(2.1), Inches(3.84), Inches(2.5),
      "What we would need",
      [[("Nothing confidential to start. The instrument runs on public "
         "models and synthetic scenarios in your domain.", GRAY)],
       [("A domain expert to audit that rewrites preserve the facts.",
         INK, True)]], size=14)

block(s, Inches(0.62), Inches(4.85), Inches(12.1), Inches(1.65),
      "The one thing worth remembering",
      [[("A text-only check can be evaded by rewording, the obvious fix does "
         "not close it, and a single robustness score will not tell you. "
         "Until the front end is hardened, ", GRAY),
        ("the safe default has to carry the risk", INK, True),
        (".", GRAY)],
       [("Everything in this deck is measured, sample sizes are stated, and "
         "the negative results are included.", GRAY)]], size=15)
notes(s, """
Closing slide. The ask is deliberately soft because the goal of this meeting is
credibility and relationship, not a signed engagement.

The middle column is the important one. Asking them three concrete questions
turns the end of the talk into a discussion rather than a pitch, and their
answers tell you whether there is real work here.

Note the third column explicitly: nothing confidential is needed to start. That
lowers the barrier to a follow-up meeting considerably, and it is true.

Leave the last block on screen during discussion.
""")

# ===========================================================================
# APPENDIX
# ===========================================================================
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.3), SW, Pt(3))
band.fill.solid()
band.fill.fore_color.rgb = BLUE
band.line.fill.background()
band.shadow.inherit = False
tf = tb(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.9))
runs(tf.paragraphs[0], [("Appendix: method and evidence", INK, True)], 34)
tf = tb(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.6))
runs(tf.paragraphs[0],
     [("For the technical session. Validation, the decision kernel, the "
       "bifactor result, limits, and a bilingual glossary.", GRAY)], 16)
notes(s, "Divider. Skip straight past this in a 30-minute slot.")

# --- A1 reliability ---------------------------------------------------------
s = new("Is the harm space reliably measurable?", 12)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.8))
para(tf, [("The first question a referee asked, and the right one. If the "
           "seven axes are not measurable to begin with, displacement along "
           "them means nothing. We tested agreement across an open, "
           "multi-family panel.", GRAY)], size=15, first=True)

rows = [["", "ICC(2,k)", "Krippendorff \u03b1"],
        ["Financial (highest)", "0.983", "0.902"],
        ["Trust", "0.958", "0.783"],
        ["Physical / autonomy (lowest)", "0.893", "about 0.57"],
        ["Overall", ("0.969", INK, True), ("0.836", INK, True)]]
table(s, Inches(0.62), Inches(2.45), Inches(7.4), rows,
      [Inches(3.6), Inches(1.9), Inches(1.9)])

block(s, Inches(8.35), Inches(2.45), Inches(4.37), Inches(2.4),
      "Panel design",
      [[("Six models, five families: Qwen3, Qwen3-small, GLM-5, GPT-OSS, "
         "MiniMax-M2, Gemma.", GRAY)],
       [("31 scenarios, disjoint from the systems under test.", GRAY)],
       [("Open managed API, fully reproducible.", GRAY)]], size=13)
block(s, Inches(0.62), Inches(4.55), Inches(12.1), Inches(1.9),
      "Reading the numbers",
      [[("Overall agreement of 0.969 and test\u2013retest of 0.96 mean the "
         "axes index a structure that different model families share. They "
         "are not one model's idiosyncrasy.", GRAY)],
       [("Physical and autonomy are the weakest, near 0.57 on the stricter "
         "statistic. We report the weakest axis rather than only the "
         "average.", GRAY)]], size=14)
notes(s, """
Engineers ask this first, so have it ready. ICC and Krippendorff alpha are
inter-rater reliability statistics: they ask whether independent judges agree
beyond chance.

The honest detail to volunteer: physical and autonomy are the weakest axes. If
someone is going to use one axis for a decision, that matters, and hiding it
would be the wrong call.

Reproducibility point: the panel runs on an open managed API rather than a
private vendor account, so the validation can be repeated by a third party.
""")

# --- A2 kernel -------------------------------------------------------------
s = new("The exploit survives a principled decision kernel", 13)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.75))
para(tf, [("A natural objection: perhaps the problem is the scoring, and a "
           "properly structured decision procedure would be immune. We "
           "routed the same scenarios through a real typed decision kernel "
           "and re-measured.", GRAY)], size=15, first=True)

rows = [["Measure", "Result"],
        ["Verdicts moving from forbid or avoid to permissive",
         ("13.7 percent", ORANGE, True)],
        ["Items rated forbid, before and after softening",
         ("44 \u2192 27", ORANGE, True)],
        ["Scenarios evaluated", "n = 161"]]
table(s, Inches(0.62), Inches(2.4), Inches(9.0), rows,
      [Inches(6.0), Inches(3.0)])

block(s, Inches(0.62), Inches(4.15), Inches(12.1), Inches(2.3),
      "Why this matters more than the score result",
      [[("The scalar score moving is a measurement artefact you could argue "
         "about. A ", GRAY), ("typed verdict", INK, True),
        (" moving from forbid to permissive is a decision changing.", GRAY)],
       [("It also localises the fault. The kernel's logic is sound; it "
         "inherits a movable front end. Structure downstream does not "
         "protect you from manipulation upstream, which is why we name the "
         "fact-extraction stage as the hardening target.", GRAY)]], size=14)
notes(s, """
This is the answer to "surely a proper rule engine fixes this". No: the rule
engine is fine and still gets the wrong answer, because it is fed a moved
representation.

The transferable engineering lesson is about interface contracts. If a
downstream component assumes its input is presentation-invariant and the
upstream one does not guarantee that, the assumption is unenforced.

Worth saying: this is the same failure pattern as a sound controller acting on a
miscalibrated sensor.
""")

# --- A3 bifactor -----------------------------------------------------------
s = new("Reliable does not mean distinct", 14)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.75))
para(tf, [("High agreement says the axes are measurable. It does not say "
           "they are independent. We registered a test for that, and it "
           "changed what we claim.", GRAY)], size=15, first=True)

rows = [["Finding", "Value"],
        ["A single general dimension separates the cases (AUROC)", "0.856"],
        ["Named axes predictable from that one dimension alone",
         ("five axes at 0.98 or above", ORANGE, True)],
        ["Variance carried by the first component in the panel's own matrix",
         ("54.2 percent", ORANGE, True)],
        ["Axis that stays genuinely specific",
         ("physical", TEAL, True)]]
table(s, Inches(0.62), Inches(2.4), Inches(9.6), rows,
      [Inches(6.6), Inches(3.0)])

block(s, Inches(0.62), Inches(4.35), Inches(12.1), Inches(2.1),
      "What we did with an unwelcome result",
      [[("Most of the apparent multi-dimensionality is one general "
         "valence channel. That weakens a claim we would have liked to keep, "
         "so we demoted it in print rather than around it.", GRAY)],
       [("It also sharpens the practical advice: a single score is mostly "
         "that one channel, which is ", GRAY),
        ("why", INK, True),
        (" one score cannot see the specific vulnerabilities. And physical "
         "harm remains separately informative.", GRAY)]], size=14)
notes(s, """
Present this as a result that went against you and was published anyway. In a
Japanese engineering culture that values not overselling, this slide does more
for credibility than any positive number in the deck.

The technical content: a bifactor structure means one general factor plus
axis-specific residuals. Here the general factor dominates and most residuals
are near zero, so the axes are reliable but largely redundant.

Note that this strengthens rather than weakens the procurement argument on slide
8, and say so.
""")

# --- A4 limits -------------------------------------------------------------
s = new("What we have not shown", 15)
tf = tb(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.5))
para(tf, [("The four claims that bound everything above.", GRAY)],
     size=15, first=True)
tf = tb(s, Inches(0.62), Inches(2.05), Inches(12.1), Inches(4.4))
para(tf, [("Adversarial defence is unproven at scale. ", INK, True),
          ("Averaging halves drift on natural rewording. Against softened "
           "register the audited set shows it does not hold, 3 of 3 still "
           "flip. The at-scale adversarial run is designed and not yet run.",
           GRAY)], size=15, first=True, bullet=True)
para(tf, [("Two claims rest on small samples. ", INK, True),
          ("The end-to-end threshold result is six hand-audited items. The "
           "zero-sycophancy observation is 0 of 9. Both are directionally "
           "clear and neither is a population rate.", GRAY)],
     size=15, bullet=True)
para(tf, [("Two different implementations. ", INK, True),
          ("Attack results come from model judges; the defence and bifactor "
           "results from trained per-axis scorers. Running all five tracks "
           "through one pipeline is designed, not done.", GRAY)],
     size=15, bullet=True)
para(tf, [("Corpus scope. ", INK, True),
          ("Moderation and social-norm material, predominantly English. "
           "Cross-cultural and cross-domain transfer is untested, which is a "
           "real limitation for a Japanese deployment.", GRAY)],
     size=15, bullet=True)
footnote(s, [("Every number in this deck is hedged where it appears. This "
              "slide collects the four bounds in one place.", GRAY)])
notes(s, """
Do not rush this slide and do not apologise for it. Offering the limits before
being asked is the single most effective credibility move available.

The fourth point is the one to emphasise in Kyoto: the corpora are largely
English and Western-normed. Cross-cultural transfer is untested. If Omron cares
about Japanese-language deployment, that gap is precisely a place where a
collaboration would produce something neither side has.

If asked which limitation worries you most, answer the first one honestly.
""")

# --- A5 reproducibility ----------------------------------------------------
s = new("Reproducibility and what exists today", 16)
rows = [["Component", "State"],
        ["Benchmark and five measurement tracks",
         ("published, tagged release", TEAL, True)],
        ["Validation panel on an open managed API",
         ("reproducible by a third party", TEAL, True)],
        ["Per-axis trained scorers", ("research code", INK, False)],
        ["Typed decision kernel", ("research library", INK, False)],
        ["Native re-run of all tracks in one pipeline",
         ("designed, not run", ORANGE, True)],
        ["Register-crossing generator for adversarial defence",
         ("the open engineering problem", ORANGE, True)]]
table(s, Inches(0.62), Inches(1.5), Inches(11.0), rows,
      [Inches(7.0), Inches(4.0)])

block(s, Inches(0.62), Inches(4.5), Inches(11.0), Inches(1.95),
      "Total cost of the published study",
      [[("About 8,000 model calls, under 50 US dollars per day. The point is "
         "that this class of testing is cheap enough to run routinely rather "
         "than once at certification.", GRAY)],
       [("That is the argument for making it part of a development process "
         "instead of an audit event.", INK, True)]], size=14)
notes(s, """
For engineers and for anyone assessing whether this is real. Be precise about
what is production-quality versus research code; overstating here would
undermine the honesty established earlier.

The cost line has an executive implication worth stating: at this price the
testing belongs in continuous integration, not in an annual audit.

The last two rows are the open problems. If a collaboration comes out of this
meeting, it is most likely one of those two.
""")

# --- A6 glossary -----------------------------------------------------------
s = new("Glossary of key terms (EN / JA)", 17)
tf = tb(s, Inches(0.62), Inches(1.4), Inches(12.1), Inches(0.45))
para(tf, [("Draft translations for reference. ", GRAY),
          ("Please have a native speaker verify before use.", ORANGE, True)],
     size=13, first=True)
rows = [["English", "\u65e5\u672c\u8a9e", "Meaning here"],
        ["harm space", "\u5371\u5bb3\u7a7a\u9593",
         "the seven-axis space judgments are mapped into"],
        ["systematic fault", "\u7cfb\u7d71\u7684\u6545\u969c",
         "reproducible and directional, not diluted by volume"],
        ["random fault", "\u30e9\u30f3\u30c0\u30e0\u6545\u969c",
         "the kind statistics over many cases does catch"],
        ["diagnostic coverage", "\u8a3a\u65ad\u30ab\u30d0\u30ec\u30c3\u30b8",
         "the fraction of faults a check can actually detect"],
        ["fail-safe / safe state",
         "\u30d5\u30a7\u30fc\u30eb\u30bb\u30fc\u30d5\uff0f\u5b89\u5168\u72b6\u614b",
         "escalate rather than pass when the check is untrustworthy"],
        ["threshold evasion", "\u95be\u5024\u56de\u907f",
         "rewording until the score falls below the flag threshold"],
        ["euphemism", "\u5a49\u66f2\u8868\u73fe",
         "softened wording that preserves the facts"],
        ["register (of language)", "\u6587\u4f53\u30fb\u8a9e\u8abf",
         "the style level; the actual attack vector here"],
        ["invariance", "\u4e0d\u5909\u6027",
         "the property a sound evaluator should have"],
        ["displacement", "\u5909\u4f4d",
         "how far a judgment moved when the wording changed"],
        ["inter-rater reliability",
         "\u8a55\u4fa1\u8005\u9593\u4fe1\u983c\u6027",
         "whether independent judges agree beyond chance"],
        ["calibration error",
         "\u30ad\u30e3\u30ea\u30d6\u30ec\u30fc\u30b7\u30e7\u30f3\u8aa4\u5dee",
         "gap between stated confidence and actual accuracy"],
        ["salience", "\u9855\u8457\u6027",
         "how noticeable a feature is made"]]
table(s, Inches(0.62), Inches(1.95), Inches(12.1), rows,
      [Inches(3.1), Inches(3.1), Inches(5.9)], size=12)
notes(s, """
Hand this out or leave it as a reference slide rather than presenting it.

The safety terms in the middle rows are the standardised Japanese forms used in
functional-safety standards work, which is why they are the right bridge for an
Omron audience. Even so, have a native speaker check the whole column before the
meeting; the linguistic terms in particular have several acceptable renderings
and I have chosen the ones closest to the sense used here.

If an Omron engineer corrects one of these, that is a good sign and worth
thanking them for.
""")

prs.save(OUT)
print(f"wrote {OUT}")
print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
